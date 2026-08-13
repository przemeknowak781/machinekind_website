"""Składa PDF i PPTX z gotowych PNG-ów w export/png/.

Uruchamia to `build.mjs` po zrzutach, ale skrypt działa też sam:

    python3 build-pack.py

Wymaga `pillow` i `python-pptx`.
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).parent
PNG_DIR = HERE / "export" / "png"
EXPORT = HERE / "export"

# Kadr prezentacji: 16:9 na 13,333 × 7,5 cala, czyli tyle, ile domyślnie
# ma slajd panoramiczny w PowerPoincie i Keynote.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

pngs = sorted(PNG_DIR.glob("*.png"))
if not pngs:
    raise SystemExit("Brak PNG-ów w export/png — uruchom najpierw `node build.mjs`.")

# ---- PDF: jedna plansza na stronę, w kadrze 1920 × 1080 punktów ----
# Strony liczy się w punktach (1/72 cala), więc zrzut w 2× schodzi do 1920
# szerokości logicznej i zostaje ostry przy każdym powiększeniu.
pages = []
for path in pngs:
    im = Image.open(path).convert("RGB")
    if im.width != 1920:
        im = im.resize((1920, round(im.height * 1920 / im.width)), Image.LANCZOS)
    pages.append(im)

pdf_path = EXPORT / "machinekind-plansze.pdf"
pages[0].save(
    pdf_path,
    "PDF",
    save_all=True,
    append_images=pages[1:],
    resolution=72.0,
)
print(f"PDF   → {pdf_path} ({pdf_path.stat().st_size / 1e6:.2f} MB)")

# ---- PPTX: każda plansza jako obraz na pełnym kadrze ----
deck = Presentation()
deck.slide_width = SLIDE_W
deck.slide_height = SLIDE_H
blank = deck.slide_layouts[6]

for path in pngs:
    slide = deck.slides.add_slide(blank)
    slide.shapes.add_picture(str(path), 0, 0, width=SLIDE_W, height=SLIDE_H)

pptx_path = EXPORT / "machinekind-plansze.pptx"
deck.save(pptx_path)
print(f"PPTX  → {pptx_path} ({pptx_path.stat().st_size / 1e6:.2f} MB)")
